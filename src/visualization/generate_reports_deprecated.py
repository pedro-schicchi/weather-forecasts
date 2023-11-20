import os
from pandas import read_excel
from pptx import Presentation

# Based on the Master Layout: position 0 == idx 11 and position 1 == idx 12
PLACE_HOLDER_IDX = 11

class GenerateReports:
    def __init__(self, ppt_file, image_data, fig_dir):
        self.original_ppt = ppt_file
        self.final_ppt = ppt_file.rsplit('.', 1)[0] + '_updated.pptx'
        self.image_data = image_data
        self.fig_dir = fig_dir
        
    def update_ppt(self, save=True):
        # open the PowerPoint presentation
        presentation = Presentation(self.original_ppt)
        
        # for each image in passed df
        for i, img in self.image_data.iterrows():
            # Get the desired slide (slide_number - 1 because Python uses 0-based indexing)
            slide = presentation.slides[img.slide-1]
            
            # Get the desired placeholder 
            placeholder = slide.placeholders[PLACE_HOLDER_IDX + img.position]
            
            # Insert desired picture
            fig_path = os.path.join(self.fig_dir, img.filename)
            picture = placeholder.insert_picture(fig_path)
            
            # ensure picture is not cropped
            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0
            
        if save:
            presentation.save(self.final_ppt)

if __name__ == '__main__':
    BASE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', '..')
    REF_DIR = os.path.join(BASE_DIR, 'references')
    REP_DIR = os.path.join(BASE_DIR, 'reports')
    FIG_DIR = os.path.join(REP_DIR, 'figures', 'gefs_fcst')
  
    # get list of images and where to put them
    images = read_excel(os.path.join(REF_DIR, 'manual_inputs.xlsx'), sheet_name='images')
    
    # ppt file name
    ppt_file = os.path.join(REP_DIR,'weather_report_pt.pptx')
    
    # 
    fr = GenerateReports(ppt_file, images)
    fr.update_ppt()
    
    # # open the PowerPoint presentation
    # presentation = Presentation(ppt_file)
    
    # # add images
    # for i, img_meta in images.iterrows():
    #     # Get the desired slide (slide_number - 1 because Python uses 0-based indexing)
    #     slide = presentation.slides[img_meta.slide-1]
        
    #     # Get the desired placeholder (position 0 == idx 11 and position 1 == idx 12)
    #     placeholder = slide.placeholders[11+img_meta.position]
        
    #     # Insert desired picture
    #     fig_path = os.path.join(FIG_DIR,img_meta.filename)
    #     picture = placeholder.insert_picture(fig_path)
    
    # # save presentation
    
    # presentation.save(os.path.join(REP_DIR,'test.pptx'))    